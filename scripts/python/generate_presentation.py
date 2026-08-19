import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from odf.opendocument import OpenDocumentPresentation
from odf.draw import Page, Frame, TextBox
from odf.text import P, Span
from odf.style import Style, MasterPage, PageLayout, PageLayoutProperties, GraphicProperties, TextProperties, ParagraphProperties

def build_pptx():
    os.makedirs("docs/presentation", exist_ok=True)
    pptx_path = "docs/presentation/unpd_defense_presentation.pptx"
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Strict Academic Palette (Deep Oxford Navy / Slate / Crisp Silver / Crimson Accent)
    BG_DARK = RGBColor(15, 23, 42)         # #0F172A (Deep Slate Navy)
    CARD_BG = RGBColor(30, 41, 59)         # #1E293B (Dark Slate Card)
    CARD_BORDER = RGBColor(51, 65, 85)     # #334155 (Subtle Divider)
    TEXT_MAIN = RGBColor(248, 250, 252)    # #F8FAFC (High-Contrast White)
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 (Refined Slate Gray)
    TEXT_CODE = RGBColor(56, 189, 248)     # #38BDF8 (Sky Blue for technical identifiers)
    ACCENT_TEAL = RGBColor(45, 212, 191)   # #2DD4BF (Academic Teal Accent)
    ACCENT_BLUE = RGBColor(96, 165, 250)   # #60A5FA (Academic Slate Blue)
    ACCENT_AMBER = RGBColor(251, 191, 36)  # #FBBF24 (Amber Metric Highlight)

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, section_num=""):
        # Section Category / Numbering
        sec_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.25))
        tf_s = sec_box.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = f"РАЗДЕЛ {section_num}".upper() if section_num else "НАУЧНО-КВАЛИФИКАЦИОННЫЙ ДОКЛАД"
        p_s.font.name = "Calibri"
        p_s.font.size = Pt(10)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_TEAL
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.55))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN

    def add_card(slide, left, top, width, height, title=None, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.0)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.15), Inches(width - 0.5), Inches(0.35))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
        return card

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)
    
    add_card(slide1, 1.0, 1.2, 11.333, 5.1, border_color=CARD_BORDER)
    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.333), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "НАУЧНО-ИНЖЕНЕРНАЯ ДИССЕРТАЦИОННАЯ ПРЕЗЕНТАЦИЯ"
    p0.font.name = "Calibri"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_TEAL
    p0.space_after = Pt(14)
    
    p1 = tf.add_paragraph()
    p1.text = "Архитектура высоконадежного драйвера ядра Windows x64\nна базе стандарта C++20 с полиморфным диспетчером памяти\nи 4-уровневой моделью трансляции MMU"
    p1.font.name = "Calibri"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.space_after = Pt(18)
    
    p2 = tf.add_paragraph()
    p2.text = "Проект: Unsolicited Non-Paged Driver (UNPD)  |  Специализация: Системное программирование и отказоустойчивые ОС"
    p2.font.name = "Calibri"
    p2.font.size = Pt(13)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "Докладчик: EvilEmployer  |  Репозиторий: LucPrusPPi/unsolicited-non-paged-driver  |  Ветка: rework/production-kernel-v2"
    p3.font.name = "Calibri"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: Problem Statement & Objectives
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_header(slide2, "Постановка проблемы и научные цели исследования", "1")
    
    add_card(slide2, 0.8, 1.4, 5.6, 5.4, title="Проблематика традиционной разработки ядра")
    tb_p = slide2.shapes.add_textbox(Inches(1.05), Inches(1.95), Inches(5.1), Inches(4.6))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    
    problems = [
        ("Небезопасность процедурного C-подхода", "Отсутствие инкапсуляции и деструкторов в традиционном коде ядра влечет риск утечек ресурсов (Pool Leaks) и повторного освобождения (Double Free)."),
        ("Ограничения среды исполнения Ring-0", "Аппаратный запрет механизма C++ исключений и RTTI в режиме ядра вынуждает применять небезопасные C-паттерны при обработке системных ошибок."),
        ("Высокая латентность IPC и переключения контекста", "Классический ввод-вывод IRP требует постоянного копирования буферов и смены уровней привилегий, снижая пропускную способность потока данных."),
        ("Сложность автономной верификации без гипервизора", "Отсутствие эмуляции таблиц страниц x86-64 в юзермоде делает модульное тестирование в CI невозможным без физической ВМ и отладчика WinDbg.")
    ]
    for idx, (title, desc) in enumerate(problems):
        p = tf_p.paragraphs[0] if idx == 0 else tf_p.add_paragraph()
        p.text = f"1.{idx+1}. {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        
    add_card(slide2, 6.8, 1.4, 5.7, 5.4, title="Научно-инженерные решения проекта UNPD")
    tb_s = slide2.shapes.add_textbox(Inches(7.05), Inches(1.95), Inches(5.2), Inches(4.6))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    
    solutions = [
        ("Freestanding ядро C++20 (kstd)", "Реализация легковесных RAII-контейнеров, умных указателей и монады expected без зависимостей от стандартного C-Runtime ядра."),
        ("Полиморфная стратегия памяти (IMemoryEngine)", "Объектно-ориентированная абстракция над MDL Zero-Copy, 4-классовыми Lookaside-кэшами и динамически отслеживаемым NonPagedPoolNx."),
        ("Аппаратно-точный 4-уровневый разбор MMU x86-64", "Алгоритм трансляции CR3 PML4 -> PDPTE -> PDE -> PTE с поддержкой супер-страниц 1GB/2MB и побайтовым чанкингом границ страниц."),
        ("Изолированный 64MB эмулятор Virtual MMU Sandbox", "Модульный стенд в юзермоде для детерминированного тестирования трансляции страниц и аппаратных прерываний #PF в GitHub CI.")
    ]
    for idx, (title, desc) in enumerate(solutions):
        p = tf_s.paragraphs[0] if idx == 0 else tf_s.add_paragraph()
        p.text = f"2.{idx+1}. {title}: "
        p.font.name = "Calibri"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: System Topology
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_header(slide3, "Сквозная архитектурная топология системы", "2")
    
    add_card(slide3, 0.8, 1.4, 11.7, 1.6, title="Уровень пользовательского режима (User-Mode Client SDK)")
    tb = slide3.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(11.2), Inches(0.95))
    p = tb.text_frame.paragraphs[0]
    p.text = "• Класс DriverClient: 16 типизированных методов (чтение/запись по CR3, постановка KAPC, очистка PiDDB, выделение MDL, слэб-кэши, пинг)\n• Класс SharedRingSession: Высокоуровневый RAII-контейнер lockless кольцевого канала с поддержкой Mock Loopback для автономного CI"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN

    add_card(slide3, 0.8, 3.2, 11.7, 2.1, title="Уровень ядра Windows (Ring-0 Kernel Dispatch & Engines)")
    tb = slide3.shapes.add_textbox(Inches(1.05), Inches(3.7), Inches(11.2), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "• Диспетчер IRP (16 кодов IOCTL 0x800..0x80F): Полная изоляция в SEH (__try/__except), валидация ProbeForRead/Write, 16MB лимит\n• Полиморфный диспетчер памяти IMemoryEngine: MdlMemoryEngine (Zero-Copy), SlabMemoryEngine (Lookaside), PoolMemoryEngine\n• Подсистема MMU Paging & Cr3Walker: Прямой 4-уровневый разбор таблиц страниц PML4 -> PDPTE -> PDE -> PTE без прикрепления\n• Подсистемы KernelApc и StealthCleaners: Асинхронные APC с Rundown-защитой, AVL-ребалансировка PiDDBCacheTable"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN

    add_card(slide3, 0.8, 5.5, 11.7, 1.4, title="Аппаратно-зависимый уровень (MASM64 / Hardware Primitives)")
    tb = slide3.shapes.add_textbox(Inches(1.05), Inches(6.0), Inches(11.2), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = "• Инвалидация кэшей TLB (invlpg, wbinvd, перезагрузка CR3)  |  Аппаратный SSE4.2 CRC32 (crc32 rax, rdx)  |  Барьеры памяти (mfence, lfence, sfence)\n• Прямой доступ к регистрам процессора: CR0..CR8, DR0..DR7, XCR0, IA32_EFER, IDTR, GDTR, TR"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 4: Freestanding Core
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_header(slide4, "Freestanding библиотека компонентов ядра (unpd::kstd)", "3")
    
    add_card(slide4, 0.8, 1.4, 5.6, 2.6, title="1. kstd::span<T> — Непрерывный срез памяти")
    tb = slide4.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Легковесный срез непрерывной памяти без владения и без CRT\n• Полная поддержка итераторов (begin, end, rbegin, rend)\n• Методы безопасного деления (subspan, first, last) с проверкой границ\n• Полная совместимость с ядерными буферами MDL и пулом"
    tb.text_frame.paragraphs[0].font.size = Pt(10.5)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide4, 6.8, 1.4, 5.7, 2.6, title="2. kstd::expected<T, NTSTATUS> — Монада ошибок")
    tb = slide4.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Монадический контейнер значения или кода статуса NTSTATUS\n• Полная замена механизма C++ исключений, запрещенного в Ring-0\n• Методы has_value(), value(), error(), value_or(fallback)\n• Специализация kstd::expected<void, E> для статусных процедур"
    tb.text_frame.paragraphs[0].font.size = Pt(10.5)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide4, 0.8, 4.2, 5.6, 2.6, title="3. kstd::unique_ptr<T, Tag> — RAII пула ядра")
    tb = slide4.shapes.add_textbox(Inches(1.05), Inches(4.7), Inches(5.1), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Инкапсулирует владение динамическим блоком невыгружаемого пула\n• Автоматический вызов ExFreePoolWithTag при выходе из скоупа\n• Семантика чистого перемещения (Move-only), запрет копирования\n• Исключает утечки пула при ранних возвратах из функций"
    tb.text_frame.paragraphs[0].font.size = Pt(10.5)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide4, 6.8, 4.2, 5.7, 2.6, title="4. Концепты времени компиляции (C++20 Concepts)")
    tb = slide4.shapes.add_textbox(Inches(7.05), Inches(4.7), Inches(5.2), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• integral<T>: Ограничение шаблонных целочисленных аргументов\n• pointer<T>: Статическая валидация указательных типов данных\n• same_as<T, U>: Строгая эквивалентность типов в шаблонах\n• trivially_copyable<T>: Гарантия безопасности memcpy для структур\n• invocable<F, Args...>: Проверка сигнатур callback-функций ядра"
    tb.text_frame.paragraphs[0].font.size = Pt(10.5)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 5: MMU & CR3 Walker
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_header(slide5, "Математическая модель трансляции MMU x86-64 и CR3 Walker", "4")
    
    add_card(slide5, 0.8, 1.4, 11.7, 5.4, title="Спецификация 4-уровневой трансляции канонического адреса")
    tb = slide5.shapes.add_textbox(Inches(1.05), Inches(1.95), Inches(11.2), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Схема адресации: PML4E [47:39] -> PDPTE [38:30] -> PDE [29:21] -> PTE [20:12] -> Offset [11:0]"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(10)
    
    points = [
        ("Проверка каноничности адреса (IsCanonical):", "Верификация знакового расширения 48-го бита на старшие 16 бит. Предотвращает немедленный аппаратный сбой #GP (General Protection Fault) при обработке невалидных адресов."),
        ("Трансляция 1GB Huge Pages (PDPTE.LargePage = 1):", "Физический адрес = (PDPTE & 0x000FFFFFC0000000) | (VA & 0x3FFFFFFF). Мгновенный расчет физического смещения без обхода нижележащих уровней PDE и PTE."),
        ("Трансляция 2MB Large Pages (PDE.LargePage = 1):", "Физический адрес = (PDE & 0x000FFFFFFFE00000) | (VA & 0x1FFFFF)."),
        ("Трансляция 4KB Standard Pages (PTE):", "Физический адрес = (PTE & 0x000FFFFFFFFFF000) | (VA & 0xFFF)."),
        ("Алгоритм Page-Boundary Clamping (Чанкинг границ):", "Размер копируемого фрагмента строго ограничен границей текущей 4KB страницы: chunk = min(remaining, 4096 - (VA & 0xFFF)). Полностью исключает выход за пределы страницы в невыделенные физические фреймы."),
        ("RAII-обертка PhysicalMemoryMapping<T>:", "Инкапсулирует пары системных вызовов MmMapIoSpace и MmUnmapIoSpace с гарантированным освобождением маппинга при любых путях выхода из функции.")
    ]
    for title, desc in points:
        p = tf.add_paragraph()
        p.text = f"• {title} "
        p.font.name = "Calibri"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: Lockless Ring Channel
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_header(slide6, "Кольцевой канал разделяемой памяти и двойная буферизация", "5")
    
    add_card(slide6, 0.8, 1.4, 5.6, 5.4, title="Архитектура кольцевого буфера и кэш-линий")
    tb = slide6.shapes.add_textbox(Inches(1.05), Inches(1.95), Inches(5.1), Inches(4.6))
    tb.text_frame.paragraphs[0].text = "• Разделение кэш-линий alignas(64):\n  Указатели RequestHead/Tail и ResponseHead/Tail изолированы по независимым 64-байтовым кэш-линиям, что полностью устраняет деградацию производительности от ложного разделения (False Sharing) между ядрами CPU.\n\n• Аппаратные барьеры сериализации памяти:\n  - Барьер загрузки: UnpdLoadFence (_mm_lfence)\n  - Барьер сохранения: UnpdStoreFence (_mm_sfence)\n  - Атомарная смена буферов: UnpdFastSwapBarrier (mfence)\n\n• Lockless Ring Wrap-Around:\n  Индекс слота вычисляется как slot = tail % CAPACITY с контролем переполнения (head - tail) >= CAPACITY."
    tb.text_frame.paragraphs[0].font.size = Pt(10.5)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide6, 6.8, 1.4, 5.7, 5.4, title="Экспериментальные задержки (bench_latency.py)")
    tb = slide6.shapes.add_textbox(Inches(7.05), Inches(1.95), Inches(5.2), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Результаты микропрофилирования (2000 итераций, субмикросекундная точность):\n"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(8)
    
    benchmarks = [
        ("IOCTL Ping Roundtrip:", "Мин: 0.30 µs  |  Среднее: 0.39 µs  |  P95: 0.40 µs  |  P99: 0.60 µs"),
        ("Атомарная смена буфера (Swap):", "Мин: 4.30 µs  |  Среднее: 4.44 µs  |  P95: 4.60 µs  |  P99: 5.10 µs"),
        ("Маппинг памяти в ядре (Map):", "Мин: 0.40 µs  |  Среднее: 0.47 µs  |  P95: 0.50 µs  |  P99: 0.70 µs"),
        ("Пропускная способность:", "Механизм Zero-Copy обеспечивает обмен пакетами на субмикросекундных задержках без переключения IRQL и контекста прерываний.")
    ]
    for title, desc in benchmarks:
        p = tf.add_paragraph()
        p.text = f"• {title}\n  "
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_CODE

    # =========================================================================
    # SLIDE 7: Polymorphic Memory Strategy
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_header(slide7, "Полиморфная стратегия управления памятью (IMemoryEngine)", "6")
    
    add_card(slide7, 0.8, 1.4, 5.6, 2.6, title="1. MdlMemoryEngine (Zero-Copy MDL)")
    tb = slide7.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.1), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Выделение физических страниц через MmAllocatePagesForMdlEx\n• Отображение в процесс через MmMapLockedPagesSpecifyCache\n• Флаг MdlMappingNoExecute (аппаратная защита DEP)\n• Безопасный анмаппинг через ProcessAttachmentGuard"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide7, 6.8, 1.4, 5.7, 2.6, title="2. SlabMemoryEngine (Lookaside кэши)")
    tb = slide7.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• O(1) Lookaside-списки (NPAGED_LOOKASIDE_LIST)\n• 4 Класса фиксированных блоков: 64B, 256B, 1024B, 4096B\n• Выделенные тэги пула: '1LSU', '2LSU', '3LSU', '4LSU'\n• Минимальная фрагментация памяти ядра"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide7, 0.8, 4.2, 5.6, 2.6, title="3. PoolMemoryEngine (Tracked NonPaged)")
    tb = slide7.shapes.add_textbox(Inches(1.05), Inches(4.7), Inches(5.1), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Выделение NonPagedPoolNx через ExAllocatePool2\n• Таблица 64-битных дескрипторов аллокаций\n• Отслеживание активных аллокаций и глобальная телеметрия\n• Полная защита спинлоками KeAcquireSpinLock"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide7, 6.8, 4.2, 5.7, 2.6, title="4. DirectNeitherEngine (SEH Probed)")
    tb = slide7.shapes.add_textbox(Inches(7.05), Inches(4.7), Inches(5.2), Inches(1.9))
    tb.text_frame.paragraphs[0].text = "• Обработка буферов METHOD_NEITHER и METHOD_IN_DIRECT\n• Проверка указателей ProbeForRead / ProbeForWrite\n• Изоляция всех обращений в блоках __try / __except\n• Исключение системных сбоев при краше юзермода"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 8: Zero-BSOD Invariants
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_header(slide8, "Доказательство отказоустойчивости ядра (Zero-BSOD Invariants)", "7")
    
    add_card(slide8, 0.8, 1.4, 11.7, 5.4, title="Матрица нейтрализации системных кодов BugCheck")
    tb = slide8.shapes.add_textbox(Inches(1.05), Inches(1.95), Inches(11.2), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    
    invariants = [
        ("0x0A / 0xD1 (IRQL_NOT_LESS_OR_EQUAL):", "Все пулы памяти и структуры данных строго аллоцируются с флагом POOL_FLAG_NON_PAGED. Код, исполняемый на DISPATCH_LEVEL, гарантированно не обращается к выгружаемой (Paged) памяти."),
        ("0x1E (KMODE_EXCEPTION_NOT_HANDLED):", "Все операции с пользовательскими указателями и чтение дескрипторов защищены Structured Exception Handling (__try / __except (EXCEPTION_EXECUTE_HANDLER))."),
        ("0x3B (SYSTEM_SERVICE_EXCEPTION):", "Валидация выравнивания адресов и обязательный вызов ProbeForRead / ProbeForWrite перед любым доступом к пользовательским буферам."),
        ("0x50 (PAGE_FAULT_IN_NONPAGED_AREA):", "Проверка каноничности адреса IsCanonical() и побайтовое ограничение размера каждого трансфера до конца текущей 4KB-страницы (Page-Boundary Clamping)."),
        ("0x7E (SYSTEM_THREAD_EXCEPTION):", "Обработчик KernelApcRundown гарантирует вызов ExFreePoolWithTag и предотвращение утечек памяти при преждевременном завершении потока до доставки APC."),
        ("0x109 (CRITICAL_STRUCTURE_CORRUPTION):", "Отсутствие статических перехватов (Zero Static Hooks); все манипуляции используют нативные AVL-таблицы ядра.")
    ]
    for idx, (title, desc) in enumerate(invariants):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {title} "
        p.font.name = "Calibri"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: Verification & Tests
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9)
    add_header(slide9, "Результаты экспериментальной верификации", "8")
    
    add_card(slide9, 0.8, 1.4, 3.7, 5.4, title="GoogleTest (73/73 Успешно)")
    tb = slide9.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(3.3), Inches(4.6))
    tb.text_frame.paragraphs[0].text = "100% Успешно (11 сьютов):\n• IoctlTest (6/6)\n• PageEngineTest (7/7)\n• FuzzingTest (6/6)\n• StressTest (6/6 — 16 потоков)\n• MmuPagingTest (13/13)\n• StealthTest (4/4)\n• MmuAdvancedTest (2/2)\n• VirtualMmuTest (12/12)\n• KstdTest (5/5)\n• SharedMemoryChannelTest (6/6)\n• ClientIntegrationTest (6/6)"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide9, 4.8, 1.4, 3.7, 5.4, title="Python & Fuzzing (8/8 Успешно)")
    tb = slide9.shapes.add_textbox(Inches(5.0), Inches(1.95), Inches(3.3), Inches(4.6))
    tb.text_frame.paragraphs[0].text = "Автоматизация и аудит:\n• test_python_tools.py (8/8)\n• fuzz_runner.py:\n  65 граничных векторов +\n  1000 случайных мутаций (0 сбоев)\n• verify_pe.py:\n  - Native Driver Subsystem\n  - ASLR (DynamicBase): True\n  - DEP/NX (NXCompat): True\n  - Control Flow Guard: True\n  - Checksum: 0x1167d (Valid)"
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    add_card(slide9, 8.8, 1.4, 3.7, 5.4, title="GitHub CI Matrix (6/6 Зеленый)")
    tb = slide9.shapes.add_textbox(Inches(9.0), Inches(1.95), Inches(3.3), Inches(4.6))
    tb.text_frame.paragraphs[0].text = "Пайплайн GitHub Actions:\n• MSVC Release: SUCCESS\n• MSVC Debug: SUCCESS\n• Clang-CL Release: SUCCESS\n• Clang-CL Debug: SUCCESS\n• Python Automation: SUCCESS\n• Code Integrity: SUCCESS\n\nКриптографическая подпись:\n100% коммитов ветки подписаны Ed25519 ключом со статусом GitHub Verified."
    tb.text_frame.paragraphs[0].font.size = Pt(10)
    tb.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 10: Conclusion
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_background(slide10)
    add_header(slide10, "Заключение и выводы квалификационной работы", "9")
    
    add_card(slide10, 1.0, 1.4, 11.333, 5.4, border_color=CARD_BORDER)
    tb = slide10.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(10.5), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ОСНОВНЫЕ НАУЧНО-ПРАКТИЧЕСКИЕ РЕЗУЛЬТАТЫ"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(12)
    
    conclusions = [
        ("1. Академическая выверенность и стандартизация:", "Разработан строгий C++20 шаблон драйвера ядра Windows x64 с полиморфной архитектурой памяти, пригодный для использования в качестве эталонного фреймворка в академических и исследовательских целях."),
        ("2. Доказанная математическая и системная надежность:", "Полная SEH-изоляция и валидация границ исключают возникновение системных сбоев (BSOD) при любых некорректных данных со стороны пользовательского режима."),
        ("3. Полная воспроизводимость верификации без гипервизора:", "Эмулятор Virtual MMU Sandbox и Mock Loopback в клиентском SDK позволяют проводить полный цикл тестирования разбора страниц и протоколов в CI без физических ВМ."),
        ("4. Высокая производительность межмодульного взаимодействия:", "Субмикросекундный отклик (0.39 µs) и lockless ring-буфер с барьерами MASM64 обеспечивают обработку потока данных в реальном времени.")
    ]
    for title, desc in conclusions:
        p = tf.add_paragraph()
        p.text = f"{title}\n"
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"  {desc}\n"
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        
    p_end = tf.add_paragraph()
    p_end.text = "ДОКЛАД ОКОНЧЕН. СПАСИБО ЗА ВНИМАНИЕ! ГОТОВ К ОТВЕТАМ НА ВОПРОСЫ КОМИССИИ."
    p_end.font.name = "Calibri"
    p_end.font.size = Pt(11)
    p_end.font.bold = True
    p_end.font.color.rgb = ACCENT_BLUE
    p_end.space_before = Pt(8)

    prs.save(pptx_path)
    print(f"[+] PowerPoint presentation successfully generated: {pptx_path}")

def build_odp():
    odp_path = "docs/presentation/unpd_defense_presentation.odp"
    doc = OpenDocumentPresentation()
    
    pl = PageLayout(name="PL16x9")
    pl.addElement(PageLayoutProperties(pagewidth="28cm", pageheight="15.75cm", margin="0cm"))
    doc.automaticstyles.addElement(pl)
    
    mp = MasterPage(name="Standard", pagelayoutname="PL16x9")
    doc.masterstyles.addElement(mp)
    
    st_title = Style(name="TitleStyle", family="presentation")
    st_title.addElement(TextProperties(fontname="Calibri", fontsize="20pt", fontweight="bold", color="#f8fafc"))
    st_title.addElement(ParagraphProperties(lineheight="120%"))
    doc.styles.addElement(st_title)

    st_body = Style(name="BodyStyle", family="presentation")
    st_body.addElement(TextProperties(fontname="Calibri", fontsize="10.5pt", color="#94a3b8"))
    st_body.addElement(ParagraphProperties(lineheight="130%"))
    doc.styles.addElement(st_body)

    slides_content = [
        ("НАУЧНО-ИНЖЕНЕРНАЯ ДИССЕРТАЦИОННАЯ ПРЕЗЕНТАЦИЯ", 
         "Unsolicited Non-Paged Driver (UNPD)\n\n"
         "Архитектура высоконадежного драйвера ядра Windows x64 на базе стандарта C++20 с полиморфным диспетчером памяти и 4-уровневой моделью трансляции MMU\n\n"
         "Докладчик: EvilEmployer\n"
         "Ветка: rework/production-kernel-v2\n"
         "Статус проекта: Production Ready (Verified)"),
        
        ("1. Постановка проблемы и научные цели исследования",
         "ПРОБЛЕМАТИКА ТРАДИЦИОННОЙ РАЗРАБОТКИ RING-0:\n"
         "1. Небезопасность процедурного C-подхода: Отсутствие инкапсуляции и деструкторов влечет риск утечек ресурсов (Pool Leaks) и UAF.\n"
         "2. Ограничения среды исполнения: Аппаратный запрет C++ исключений и RTTI в ядре вынуждает применять небезопасные C-паттерны.\n"
         "3. Высокая латентность IPC: Классический ввод-вывод IRP требует постоянного копирования буферов и смены уровней привилегий.\n"
         "4. Сложность автономной верификации: Отсутствие эмуляции страниц в юзермоде делает модульное тестирование в CI невозможным без ВМ.\n\n"
         "НАУЧНО-ИНЖЕНЕРНЫЕ РЕШЕНИЯ В UNPD:\n"
         "1. Freestanding ядро C++20 (kstd): Собственные реализации span, expected, unique_ptr и концептов без CRT ядра.\n"
         "2. Полиморфная память (IMemoryEngine): Абстракция над MDL Zero-Copy, Lookaside-кэшами, пулом и Neither I/O.\n"
         "3. 4-Уровневый разбор MMU x86-64: CR3 Walker с поддержкой 1GB/2MB страниц и побайтовым чанкингом границ страниц.\n"
         "4. Virtual MMU 64MB Sandbox: Автономный эмулятор физической памяти и TLB с генерацией #PF в CI без гипервизора."),
        
        ("2. Сквозная архитектурная топология системы",
         "УРОВЕНЬ ПОЛЬЗОВАТЕЛЬСКОГО РЕЖИМА (include/unpd/client.hpp):\n"
         "• Класс DriverClient: 16 типизированных методов (CR3 R/W, APC, Stealth, Slabs, MDL, Stats, Ping).\n"
         "• Класс SharedRingSession: Высокоуровневый RAII-контейнер lockless кольцевого канала с Mock Loopback.\n\n"
         "УРОВЕНЬ ЯДРА WINDOWS (Ring-0 Kernel Core, src/driver/):\n"
         "• Диспетчер IRP (16 кодов IOCTL 0x800..0x80F): Полная изоляция в SEH (__try/__except) и ProbeForRead/Write.\n"
         "• Иерархия IMemoryEngine: MdlMemoryEngine (Zero-Copy), SlabMemoryEngine (Lookaside), PoolMemoryEngine.\n"
         "• MMU Paging Engine & Cr3Walker: Прямой 4-уровневый разбор PML4 -> PDPTE -> PDE -> PTE.\n"
         "• Подсистемы KernelApc и StealthCleaners: Асинхронные APC с Rundown-защитой, AVL-ребалансировка PiDDB.\n\n"
         "АППАРАТНО-ЗАВИСИМЫЙ УРОВЕНЬ (MASM64 / x86-64):\n"
         "• Инвалидация TLB (invlpg, wbinvd, CR3), аппаратный SSE4.2 CRC32 (crc32 rax, rdx), барьеры mfence, lfence, sfence."),
        
        ("3. Freestanding C++20 библиотека компонентов ядра (unpd::kstd)",
         "1. kstd::span<T>:\n"
         "• Легковесный срез непрерывной памяти без владения и без CRT ядра.\n"
         "• Полный набор итераторов и методы безопасного деления (subspan, first, last) с проверкой границ.\n\n"
         "2. kstd::expected<T, NTSTATUS>:\n"
         "• Монадический контейнер значения или кода статуса NTSTATUS вместо C++ исключений.\n"
         "• Методы has_value(), value(), error(), value_or(fallback), специализация expected<void, E>.\n\n"
         "3. kstd::unique_ptr<T, Tag>:\n"
         "• RAII умный указатель с автоматическим вызовом ExFreePoolWithTag при разрушении (Move-only).\n\n"
         "4. C++20 Concepts:\n"
         "• integral<T>, pointer<T>, same_as<T, U>, trivially_copyable<T>, invocable<F, Args...>."),
        
        ("4. Математическая модель трансляции MMU x86-64 и CR3 Walker",
         "Схема адресации: PML4E [47:39] -> PDPTE [38:30] -> PDE [29:21] -> PTE [20:12] -> Offset [11:0]\n\n"
         "• Проверка каноничности адреса (IsCanonical): Верификация 48-битного знакового расширения, исключающая сбой #GP.\n"
         "• 1GB Huge Pages (PDPTE.LargePage = 1): Физический адрес = (PDPTE & 0x000FFFFFC0000000) | (VA & 0x3FFFFFFF).\n"
         "• 2MB Large Pages (PDE.LargePage = 1): Физический адрес = (PDE & 0x000FFFFFFFE00000) | (VA & 0x1FFFFF).\n"
         "• 4KB Standard Pages: Физический адрес = (PTE & 0x000FFFFFFFFFF000) | (VA & 0xFFF).\n"
         "• Алгоритм Page-Boundary Clamping: Чанкинг строго ограничен границей 4KB страницы (min(rem, 4096 - offset)).\n"
         "• RAII PhysicalMemoryMapping<T>: Автоматическое управление парами MmMapIoSpace / MmUnmapIoSpace."),
        
        ("5. Кольцевой канал разделяемой памяти и двойная буферизация",
         "АРХИТЕКТУРА КОЛЬЦА И ВЫРАВНИВАНИЕ:\n"
         "• Разделение кэш-линий alignas(64): Указатели RequestHead/Tail и ResponseHead/Tail изолированы по 64B кэш-линиям, исключая False Sharing.\n"
         "• Аппаратные барьеры сериализации: LoadFence (lfence), StoreFence (sfence), FastSwapBarrier (mfence).\n"
         "• Lockless Ring Wrap-Around: slot = tail % CAPACITY с контролем переполнения (head - tail) >= CAPACITY.\n\n"
         "ЭКСПЕРИМЕНТАЛЬНЫЕ ЗАДЕРЖКИ (bench_latency.py, 2000 итераций):\n"
         "• IOCTL Ping Roundtrip: Мин 0.30 µs  |  Среднее 0.39 µs  |  P95: 0.40 µs  |  P99: 0.60 µs\n"
         "• Атомарная смена буфера (Swap): Мин 4.30 µs  |  Среднее 4.44 µs  |  P95: 4.60 µs  |  P99: 5.10 µs\n"
         "• Маппинг памяти в ядре (Map): Мин 0.40 µs  |  Среднее 0.47 µs  |  P95: 0.50 µs  |  P99: 0.70 µs"),
        
        ("6. Полиморфная стратегия управления памятью (IMemoryEngine)",
         "1. MdlMemoryEngine (Zero-Copy MDL):\n"
         "• Выделение физических страниц (MmAllocatePagesForMdlEx), маппинг с флагом MdlMappingNoExecute (DEP).\n"
         "• Безопасный анмаппинг в контексте процесса-владельца через ProcessAttachmentGuard.\n\n"
         "2. SlabMemoryEngine (Lookaside кэши):\n"
         "• O(1) Lookaside-списки (NPAGED_LOOKASIDE_LIST) для 4 классов: 64B, 256B, 1024B, 4096B.\n\n"
         "3. PoolMemoryEngine (Tracked NonPaged):\n"
         "• Выделение NonPagedPoolNx через ExAllocatePool2 с 64-битной таблицей дескрипторов.\n\n"
         "4. DirectNeitherEngine (SEH Probed):\n"
         "• Валидация ProbeForRead/Write с изоляцией в блоках __try / __except."),
        
        ("7. Доказательство отказоустойчивости ядра (Zero-BSOD Invariants)",
         "МАТРИЦА НЕЙТРАЛИЗАЦИИ СИСТЕМНЫХ СБОЕВ ЯДРА:\n"
         "• 0x0A / 0xD1 (IRQL_NOT_LESS_OR_EQUAL): Все пулы строго аллоцируются с POOL_FLAG_NON_PAGED. Код на DISPATCH_LEVEL гарантированно не обращается к paged памяти.\n"
         "• 0x1E (KMODE_EXCEPTION_NOT_HANDLED): Все операции защищены Structured Exception Handling (__try / __except).\n"
         "• 0x3B (SYSTEM_SERVICE_EXCEPTION): Валидация выравнивания и вызовы ProbeForRead/Write перед доступом.\n"
         "• 0x50 (PAGE_FAULT_IN_NONPAGED_AREA): Валидация каноничности IsCanonical() и чанкинг границ 4KB.\n"
         "• 0x7E (SYSTEM_THREAD_EXCEPTION): Rundown-клинер освобождает KAPC из пула при смерти потока.\n"
         "• 0x109 (CRITICAL_STRUCTURE_CORRUPTION): Отсутствие статических хуков; нативные AVL-деревья."),
        
        ("8. Результаты экспериментальной верификации",
         "GOOGLETEST СЬЮТ (73/73 УСПЕШНО, 11 СЬЮТОВ):\n"
         "• IoctlTest (6), PageEngineTest (7), FuzzingTest (6), StressTest (6 — 16 потоков),\n"
         "  MmuPagingTest (13), StealthTest (4), MmuAdvancedTest (2), VirtualMmuTest (12),\n"
         "  KstdTest (5), SharedMemoryChannelTest (6), ClientIntegrationTest (6).\n\n"
         "PYTHON TOOLING & FUZZING (8/8 УСПЕШНО):\n"
         "• fuzz_runner.py: 65 граничных векторов + 1000 случайных мутаций (0 сбоев).\n"
         "• verify_pe.py: Native Driver, ASLR, DEP/NX, Control Flow Guard, Checksum 0x1167d.\n\n"
         "GITHUB ACTIONS CI MATRIX (6/6 ЗЕЛЕНЫЙ):\n"
         "• MSVC Release/Debug, Clang-CL Release/Debug, Python Automation, Code Integrity.\n"
         "• 100% коммитов подписаны верифицированным SSH-ключом Ed25519 (GitHub Verified)."),
        
        ("9. Заключение и выводы квалификационной работы",
         "ОСНОВНЫЕ НАУЧНО-ПРАКТИЧЕСКИЕ РЕЗУЛЬТАТЫ:\n"
         "1. Академическая выверенность: Разработан строгий C++20 шаблон ядра Windows x64 с полиморфной архитектурой памяти.\n"
         "2. Доказанная надежность: Изоляция всех путей исполнения исключает возникновение BSOD при любых сбоях юзермода.\n"
         "3. 100% Воспроизводимость в CI: Virtual MMU Sandbox и Mock Loopback обеспечивают тестирование без гипервизоров.\n"
         "4. Высокая производительность: Субмикросекундный отклик (0.39 µs) и lockless ring-буфер с барьерами MASM64.\n\n"
         "ДОКЛАД ОКОНЧЕН. СПАСИБО ЗА ВНИМАНИЕ! ГОТОВ К ОТВЕТАМ НА ВОПРОСЫ КОМИССИИ.")
    ]

    for idx, (title, text) in enumerate(slides_content):
        page = Page(masterpagename="Standard", name=f"Slide_{idx+1}")
        
        h_frame = Frame(width="26cm", height="1.8cm", x="1cm", y="0.8cm")
        h_tb = TextBox()
        h_p = P(text=title)
        h_p.setAttribute("stylename", "TitleStyle")
        h_tb.addElement(h_p)
        h_frame.addElement(h_tb)
        page.addElement(h_frame)
        
        b_frame = Frame(width="26cm", height="12cm", x="1cm", y="2.8cm")
        b_tb = TextBox()
        for line in text.split("\n"):
            p = P(text=line)
            p.setAttribute("stylename", "BodyStyle")
            b_tb.addElement(p)
        b_frame.addElement(b_tb)
        page.addElement(b_frame)
        
        doc.presentation.addElement(page)

    doc.save(odp_path)
    print(f"[+] OpenDocument Presentation (ODP) successfully generated: {odp_path}")

if __name__ == "__main__":
    build_pptx()
    build_odp()
